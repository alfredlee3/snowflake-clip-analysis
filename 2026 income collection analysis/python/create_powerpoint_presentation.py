"""
PowerPoint Presentation Generator for PIE Success Rate Analysis

This script creates a PowerPoint presentation with embedded chart data (not images).
The charts are native PowerPoint chart objects backed by Excel data.

Prerequisites:
    pip install python-pptx pandas openpyxl snowflake-connector-python

Usage:
    python create_powerpoint_presentation.py
"""

import snowflake.connector
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import sys
import os

# Snowflake connection configuration
SNOWFLAKE_CONFIG = {
    'user': 'ALFRED_LEE',
    'account': 'IJ90379-MISSIONLANE',
    'authenticator': 'externalbrowser',
    'database': 'EDW_DB',
    'schema': 'PUBLIC'
}

# Output file path
OUTPUT_FILE = '/Users/Alfred.Lee/Documents/github/visualizations/PIE_Success_Rate_Analysis.pptx'

# Color scheme (RGB values)
COLORS = {
    'Stmt 18': RGBColor(231, 76, 60),      # #e74c3c - red
    'Stmt 26': RGBColor(52, 152, 219),     # #3498db - blue
    'Stmt 34': RGBColor(46, 204, 113),     # #2ecc71 - green
    'Stmt 42+': RGBColor(243, 156, 18),    # #f39c12 - orange
    'Overall Stmt 18+': RGBColor(0, 0, 0)  # #000000 - black
}


def connect_to_snowflake():
    """Connect to Snowflake database"""
    try:
        print("\n[1/5] Connecting to Snowflake...")
        conn = snowflake.connector.connect(**SNOWFLAKE_CONFIG)
        print("✓ Connected to Snowflake")
        return conn
    except Exception as e:
        print(f"❌ Error connecting to Snowflake: {e}")
        sys.exit(1)


def load_success_rate_data(conn):
    """Load PIE success rate data from Snowflake"""
    try:
        print("\n[2/5] Loading PIE success rate data...")

        # Read SQL query
        sql_file = '/Users/Alfred.Lee/Documents/github/2026 income collection analysis/sql/pie_income_collection_over_time.sql'
        with open(sql_file, 'r') as f:
            query = f.read()

        # Execute query
        cursor = conn.cursor()
        cursor.execute(query)
        df = pd.DataFrame(cursor.fetchall(), columns=[col[0] for col in cursor.description])
        cursor.close()

        print(f"✓ Loaded {len(df)} rows of data")
        return df

    except Exception as e:
        print(f"❌ Error loading data: {e}")
        sys.exit(1)


def prepare_chart_data(df):
    """Prepare data in format suitable for PowerPoint chart"""
    # Add month label (1-8 instead of 0-7)
    df['MONTH_LABEL'] = df['MONTH_OFFSET'] + 1

    # Filter to statements we want to chart
    statements = [18, 26, 34, 442, 999]  # 442 = Stmt 42+, 999 = Overall 18+
    df_chart = df[df['STATEMENT_NUMBER'].isin(statements)].copy()

    # Create pivot table
    pivot_data = df_chart.pivot(
        index='MONTH_LABEL',
        columns='STATEMENT_LABEL',
        values='SUCCESS_RATE_PCT'
    ).reset_index()

    # Ensure column order
    column_order = ['Month', 'Stmt 18', 'Stmt 26', 'Stmt 34', 'Stmt 42+', 'Overall Stmt 18+']
    pivot_data = pivot_data.rename(columns={'MONTH_LABEL': 'Month'})
    pivot_data = pivot_data[column_order]

    return pivot_data


def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "PIE Success Rate Analysis"

    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 0, 0)

    # Add subtitle
    subtitle_top = Inches(3.5)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "April 2025 Cohort - 8 Month Performance"

    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_paragraph.font.size = Pt(24)
    subtitle_paragraph.font.color.rgb = RGBColor(100, 100, 100)

    # Add date
    date_top = Inches(6.5)
    date_box = slide.shapes.add_textbox(left, date_top, width, Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "February 2026"

    date_paragraph = date_frame.paragraphs[0]
    date_paragraph.alignment = PP_ALIGN.CENTER
    date_paragraph.font.size = Pt(16)
    date_paragraph.font.color.rgb = RGBColor(150, 150, 150)


def create_chart_slide(prs, chart_data):
    """Create slide with success rate line chart"""
    print("\n[3/5] Creating success rate chart slide...")

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Overall Success Rate Over Time - April 2025 Cohort"
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True

    # Prepare chart data for PowerPoint
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = [f'Month {int(m)}' for m in chart_data['Month']]

    # Add series for each statement
    for col in ['Stmt 18', 'Stmt 26', 'Stmt 34', 'Stmt 42+', 'Overall Stmt 18+']:
        chart_data_obj.add_series(col, chart_data[col].values)

    # Add chart to slide
    x, y, cx, cy = Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_obj
    ).chart

    # Chart formatting
    chart.has_legend = True
    chart.legend.position = 2  # Bottom
    chart.legend.include_in_layout = False

    # Set value axis range (70-100)
    value_axis = chart.value_axis
    value_axis.minimum_scale = 70
    value_axis.maximum_scale = 100
    value_axis.has_title = True
    value_axis.axis_title.text_frame.text = "Cumulative Success Rate (%)"

    # Set category axis title
    category_axis = chart.category_axis
    category_axis.has_title = True
    category_axis.axis_title.text_frame.text = "Months After PIE Evaluation"

    # Apply colors to series
    series_names = ['Stmt 18', 'Stmt 26', 'Stmt 34', 'Stmt 42+', 'Overall Stmt 18+']
    for idx, series_name in enumerate(series_names):
        series = chart.series[idx]

        # Set line color
        line = series.format.line
        line.color.rgb = COLORS[series_name]
        line.width = Pt(2.5)

        # Make Overall Stmt 18+ dashed
        if series_name == 'Overall Stmt 18+':
            line.dash_style = 2  # Dash style

    print("✓ Created success rate chart")


def create_summary_slide(prs, chart_data):
    """Create summary slide with key metrics"""
    print("\n[4/5] Creating summary slide...")

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Key Findings - Month 8 (240 Days)"
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True

    # Get month 8 data
    month_8 = chart_data[chart_data['Month'] == 8]

    # Create metrics text
    metrics_text = f"""Success Rates by Statement:

• Stmt 18:  {month_8['Stmt 18'].values[0]:.1f}%  (Best Performance)
• Stmt 26:  {month_8['Stmt 26'].values[0]:.1f}%
• Stmt 34:  {month_8['Stmt 34'].values[0]:.1f}%
• Stmt 42+: {month_8['Stmt 42+'].values[0]:.1f}%  (Late-Stage Accounts)

• Overall Stmt 18+: {month_8['Overall Stmt 18+'].values[0]:.1f}%

Key Insights:
✓ Success Rate = (Approved Outright + PIE with Income) / Total Population
✓ PIE income collection patterns are stable and reliable
✓ Earlier statements show stronger performance (higher approval rates)
✓ 8-month window balances observation time with data recency
"""

    # Add metrics text box
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = metrics_text
    text_frame.word_wrap = True

    # Format text
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_before = Pt(6)

    print("✓ Created summary slide")


def create_presentation(chart_data):
    """Create complete PowerPoint presentation"""
    print("\n[5/5] Creating PowerPoint presentation...")

    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create slides
    create_title_slide(prs)
    create_chart_slide(prs, chart_data)
    create_summary_slide(prs, chart_data)

    # Save presentation
    prs.save(OUTPUT_FILE)
    print(f"✓ Presentation saved to: {OUTPUT_FILE}")


def main():
    """Main execution function"""
    print("\n" + "=" * 80)
    print("POWERPOINT PRESENTATION GENERATOR - PIE SUCCESS RATE ANALYSIS")
    print("=" * 80)

    # Connect to Snowflake
    conn = connect_to_snowflake()

    # Load data
    df = load_success_rate_data(conn)
    conn.close()

    # Prepare chart data
    chart_data = prepare_chart_data(df)

    # Create presentation
    create_presentation(chart_data)

    # Success message
    print("\n" + "=" * 80)
    print("✓ SUCCESS!")
    print("=" * 80)
    print(f"\nPowerPoint presentation created: {OUTPUT_FILE}")
    print("\nThe presentation contains:")
    print("  - Slide 1: Title slide")
    print("  - Slide 2: Success rate line chart (embedded chart with data)")
    print("  - Slide 3: Key findings and metrics summary")
    print("\nThe chart is a native PowerPoint chart object (not an image).")
    print("You can click on it to edit the data or formatting.")
    print("=" * 80 + "\n")


if __name__ == "__main__":
    main()
